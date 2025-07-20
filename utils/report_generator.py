# utils/report_generator.py
import io
import pandas as pd
from datetime import datetime
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import plotly.graph_objects as go

# --- This module contains pure, non-Streamlit backend report generation functions ---

# --- PDF Generation Engine ---

class PDF(FPDF):
    """Custom PDF class to handle headers, footers, and content generation."""
    def header(self):
        self.set_font('Helvetica', 'B', 12)
        self.cell(0, 10, 'VERITAS - CMC Data Summary Report', 0, 1, 'C')
        self.set_font('Helvetica', '', 8)
        self.cell(0, 5, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", 0, 1, 'C')
        self.ln(10)

    def footer(self):
        self.set_y(-15)
        self.set_font('Helvetica', 'I', 8)
        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')
        self.set_x(-50)
        self.cell(0, 10, 'CONFIDENTIAL', 0, 0, 'R')

    def chapter_title(self, title):
        self.set_font('Helvetica', 'B', 12)
        self.cell(0, 10, title, 0, 1, 'L')
        self.ln(5)

    def chapter_body(self, body):
        self.set_font('Helvetica', '', 10)
        self.multi_cell(0, 5, body)
        self.ln()

    def add_dataframe(self, df):
        self.set_font('Helvetica', 'B', 9)
        col_width = self.w / (len(df.columns) + 1.5)
        # Header
        for col in df.columns:
            self.cell(col_width, 7, col, 1, 0, 'C')
        self.ln()
        # Data
        self.set_font('Helvetica', '', 9)
        for _, row in df.iterrows():
            for col in df.columns:
                self.cell(col_width, 6, str(row[col]), 1, 0, 'L')
            self.ln()
        self.ln(5)

def generate_pdf_report(report_data: dict) -> bytes:
    """Creates a formatted PDF report from the provided data."""
    pdf = PDF()
    pdf.add_page()
    
    # Section 1: Report Metadata
    pdf.chapter_title(f"1.0 Summary for Study: {report_data['study_id']}")
    metadata_body = (
        f"This document summarizes key data for the selected study, generated from the VERITAS system. "
        f"The data is intended for regulatory support and internal review.\n\n"
        f"**Selected Critical Quality Attribute (CQA):** {report_data['cqa']}\n"
        f"**Number of Data Points Included:** {len(report_data['data'])}"
    )
    pdf.chapter_body(metadata_body)

    # Section 2: Summary Statistics
    pdf.chapter_title("2.0 Summary Statistics")
    summary_stats = report_data['data'][[report_data['cqa']]].describe().round(3).reset_index()
    summary_stats.columns = ['Statistic', 'Value']
    pdf.add_dataframe(summary_stats)

    # Section 3: Analyst Commentary
    pdf.chapter_title("3.0 Analyst Commentary")
    pdf.chapter_body(report_data['commentary'])
    
    # Return PDF as bytes
    return pdf.output(dest='S').encode('latin-1')


# --- PowerPoint Generation Engine ---

def _add_table_to_slide(slide, df, left, top, width, height):
    """Helper function to add a pandas DataFrame as a table to a PowerPoint slide."""
    rows, cols = df.shape
    rows += 1 # Add a row for the header
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    for i in range(cols):
        table.columns[i].width = Inches(width.inches / cols)

    # Write header and data
    for c, col_name in enumerate(df.columns):
        table.cell(0, c).text = col_name
        table.cell(0, c).text_frame.paragraphs[0].font.bold = True
    for r in range(rows - 1):
        for c, col_name in enumerate(df.columns):
            table.cell(r + 1, c).text = str(df.iloc[r, c])

def generate_ppt_report(df: pd.DataFrame, report_title: str, plot_fig: go.Figure) -> bytes:
    """Generates a PowerPoint report with data and a plot."""
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "VERITAS Automated Study Report"
    slide.placeholders[1].text = f"Report Type: {report_title}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    
    # Slide 2: Data Summary Table
    content_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Data Summary (First 10 Rows)"
    summary_df = df.head(10)[['sample_id', 'batch_id', 'Purity', 'Main Impurity']]
    _add_table_to_slide(slide, summary_df, Inches(0.5), Inches(1.5), Inches(9.0), Inches(4.0))

    # Slide 3: Dynamic Plot
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Process Analysis"
    
    # Save plot to an in-memory image stream
    img_stream = io.BytesIO()
    plot_fig.write_image(img_stream, format="png", width=800, height=450, scale=2)
    img_stream.seek(0)
    
    slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(8))

    # Save the final presentation to an in-memory buffer
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io.getvalue()
