# pages/3_Regulatory_Report_Assembler.py
import streamlit as st
import pandas as pd
import time
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from statsforecast import StatsForecast
from statsforecast.models import AutoARIMA
import plotly.graph_objects as go

from utils.auth import display_compliance_footer, get_user_role
from utils.plotters import plot_spc_chart
from utils.data_generator import create_mock_hplc_data

def add_table_to_slide(slide, df, left, top, width, height):
    """
    Helper function to add a pandas DataFrame as a table to a PowerPoint slide.
    """
    rows, cols = df.shape
    rows += 1 # Add a row for the header
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths to be evenly distributed
    for i in range(cols):
        table.columns[i].width = Inches(width.inches / cols)

    # Write header row
    for i, col_name in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = col_name
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    # Write data rows
    for r in range(rows - 1):
        for c in range(cols):
            cell = table.cell(r + 1, c)
            cell.text = str(df.iloc[r, c])
            cell.text_frame.paragraphs[0].font.size = Pt(9)

def create_report_pptx(df, report_title, data_package_name):
    """
    Generates a complete PowerPoint report with titles, data, and plots.
    This function acts as the report template engine.
    """
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "VERITAS Automated Report"
    slide.placeholders[1].text = f"Report Type: {report_title}\nData Source: {data_package_name}\nGenerated: {time.strftime('%Y-%m-%d %H:%M:%S')}"
    
    # Slide 2: Summary and Data Table
    content_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Data Summary"
    summary_df = df.head(10)[['sample_id', 'batch_id', 'analyte_concentration', 'retention_time']]
    add_table_to_slide(slide, summary_df, Inches(0.5), Inches(1.5), Inches(9.0), Inches(4.0))

    # Slide 3: Statistical Process Control (SPC) Chart
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Process Stability Analysis (SPC)"
    spc_fig = plot_spc_chart(df, 'analyte_concentration')
    # Save plot to an in-memory image stream
    spc_img_stream = io.BytesIO()
    spc_fig.write_image(spc_img_stream, format="png", width=800, height=450, scale=2)
    spc_img_stream.seek(0)
    # Add the image to the slide
    slide.shapes.add_picture(spc_img_stream, Inches(1), Inches(1.5), width=Inches(8))

    # Save the final presentation to an in-memory buffer
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io.getvalue()


# --- Page Configuration and Sidebar ---
st.set_page_config(layout="wide", page_title="Report Assembler", page_icon="🧪")
with st.sidebar:
    st.title("VERITAS")
    get_user_role()
    
# --- Page Header and Introduction ---
st.title("Module 3: Regulatory Package & Report Assembler")
st.markdown("Automate the generation of GxP-compliant study reports from the Single Source of Truth.")

with st.expander("ℹ️ SME Overview: The Value of Automated Reporting"):
    st.info("""
        - **Problem:** Manual report creation is slow, prone to copy-paste errors, and creates a significant compliance burden.
        - **Solution:** This module automates the entire process. It pulls approved, QC'd data directly from the VERITAS data store, populates pre-validated templates, and manages the e-signature workflow.
        - **Commercial Value:** This drastically reduces report generation time (from weeks to minutes), eliminates transcription errors, ensures consistency, and provides a full audit trail for every report, making regulatory inspections smoother.
    """)

# --- Step 1: Configuration ---
st.subheader("Step 1: Configure Report")
col1, col2 = st.columns(2)
with col1:
    report_template = st.selectbox("Select Report Template", ("IND Study Report", "PK Analysis Summary"))
with col2:
    data_package = st.selectbox("Select Approved QC'd Data Package", ("PK_Study_2024-A (v2.1)", "Tox_Assay_Run_05 (v1.0)"))

# --- Step 2: Generation & Draft Review ---
st.subheader("Step 2: Generate Draft & Preview")
if st.button("Assemble Draft Report"):
    with st.spinner("Assembling draft... This may take a moment."):
        report_df = create_mock_hplc_data(50)
        report_bytes = create_report_pptx(report_df, report_template, data_package)
        st.session_state['draft_report_bytes'] = report_bytes
        st.session_state['draft_df'] = report_df.head() # Save a preview for the UI
    st.success("Draft report assembled and ready for review.")

# --- Workflow continuation based on session state ---
if 'draft_report_bytes' in st.session_state:
    st.markdown("##### Draft Preview (First 5 Rows)")
    st.dataframe(st.session_state['draft_df'], hide_index=True)
    
    # --- Step 3: E-Signature Workflow ---
    st.subheader("Step 3: Review and E-Sign")
    st.warning("By signing, you attest that you have reviewed the draft report and its contents are accurate to the best of your knowledge.")
    
    sig_col1, sig_col2 = st.columns(2)
    with sig_col1:
        reviewer_name = st.text_input("Your Full Name (Electronic Signature)")
    with sig_col2:
        password = st.text_input("Password", type="password")

    if reviewer_name and password:
        st.session_state['signed'] = True
        
    # --- Step 4: Finalization and Download ---
    st.subheader("Step 4: Download Final Report")
    if st.session_state.get('signed', False):
        st.download_button(
            label="⬇️ Download Final Signed Report (.pptx)",
            data=st.session_state['draft_report_bytes'],
            file_name=f"FINAL_{report_template.replace(' ', '_')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary"
        )
        st.success("Report is signed. An entry has been made in the E-Signature Log in the Compliance Hub.")
    else:
        st.info("Please provide your E-Signature to unlock the final report.")
else:
    st.info("Click 'Assemble Draft Report' to begin the reporting workflow.")
    
# --- Global Compliance Footer ---
display_compliance_footer()
